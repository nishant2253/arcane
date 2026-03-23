"""
Arcane Pitch Deck Auto-Filler
Reads the Black Elegant template and replaces every text box with Arcane content.
Usage: /tmp/pptxenv/bin/python3 fill_pptx.py
"""

from pptx import Presentation
from pptx.util import Pt
import copy

TEMPLATE = "Black Elegant and Modern Startup Pitch Deck Presentation.pptx"
OUTPUT   = "Arcane_Pitch_Deck.pptx"

# ---------------------------------------------------------------------------
# Helper: replace text in a shape while keeping the first run's font style
# ---------------------------------------------------------------------------
def set_text(slide, shape_name, new_text):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            # Grab font properties from the first run of the first paragraph
            try:
                first_para = tf.paragraphs[0]
                first_run  = first_para.runs[0] if first_para.runs else None
                saved_font = {
                    "bold":   first_run.font.bold   if first_run else None,
                    "size":   first_run.font.size   if first_run else None,
                    "color":  first_run.font.color.rgb
                               if first_run and first_run.font.color and first_run.font.color.type
                               else None,
                }
            except Exception:
                saved_font = {"bold": None, "size": None, "color": None}

            # Clear all paragraphs
            tf.clear()
            para = tf.paragraphs[0]
            run  = para.add_run()
            run.text = new_text

            # Restore font
            if saved_font["bold"] is not None:
                run.font.bold = saved_font["bold"]
            if saved_font["size"]:
                run.font.size = saved_font["size"]
            if saved_font["color"]:
                from pptx.dml.color import RGBColor
                run.font.color.rgb = saved_font["color"]
            return
    # shape not found — silent skip


# ---------------------------------------------------------------------------
# Slide content definitions
# ---------------------------------------------------------------------------
def fill_slide_1(slide):
    """Title slide"""
    set_text(slide, "TextBox 9",  "AI TRADING PLATFORM ON HEDERA")
    set_text(slide, "TextBox 11", "Arcane")
    set_text(slide, "TextBox 12", "Presented by Nishant Gupta")
    set_text(slide, "TextBox 13", "nishantgupta1965@gmail.com")
    set_text(slide, "TextBox 14", "AI & Agents Track  |  Hedera Hackathon 2025")
    set_text(slide, "TextBox 15", "github.com/nishantgupta/arcane")


def fill_slide_2(slide):
    """Agenda"""
    set_text(slide, "TextBox 8",  "TODAY'S AGENDA")
    set_text(slide, "TextBox 9",  "What We're Building")
    set_text(slide, "TextBox 11", "The Problem")
    set_text(slide, "TextBox 13", "Our Solution")
    set_text(slide, "TextBox 15", "Hedera Services")
    set_text(slide, "TextBox 17", "Market Size")
    set_text(slide, "TextBox 19", "Competition")
    set_text(slide, "TextBox 21", "Key Differentiators")
    set_text(slide, "TextBox 23", "Traction & Validation")
    set_text(slide, "TextBox 25", "Revenue Model")
    set_text(slide, "TextBox 27", "Roadmap & Use of Funds")


def fill_slide_3(slide):
    """Introduction"""
    set_text(slide, "TextBox 10", "INTRODUCTION")
    set_text(slide, "TextBox 9",
        "Arcane is an AI-powered trading platform built natively on Hedera Hashgraph. "
        "We let anyone deploy autonomous AI trading agents — no code required — "
        "with every signal and trade recorded immutably on Hedera HCS."
    )
    set_text(slide, "TextBox 11",
        "From natural-language agent creation to an NFT strategy marketplace, "
        "Arcane makes institutional-grade quant trading accessible, transparent, "
        "and verifiable for the first time."
    )


def fill_slide_4(slide):
    """Problem Statement"""
    set_text(slide, "TextBox 12", "PROBLEM STATEMENT")
    set_text(slide, "TextBox 10", "Opacity in AI Trading")
    set_text(slide, "TextBox 11",
        "AI trading bots operate as black boxes. Traders cannot verify whether "
        "signals are genuine or cherry-picked. There is no tamper-proof audit trail."
    )
    set_text(slide, "TextBox 13", "Zero Trust in AI Agents")
    set_text(slide, "TextBox 14",
        "Centralized platforms can manipulate or falsify performance history. "
        "Users have no on-chain proof that the agent actually executed the strategy shown."
    )
    set_text(slide, "TextBox 15", "Retail Locked Out of Quant Strategies")
    set_text(slide, "TextBox 16",
        "Professional quant strategies require coding skills and expensive infrastructure. "
        "Retail traders cannot access or monetize institutional-quality algorithms."
    )


def fill_slide_5(slide):
    """Solutions"""
    set_text(slide, "TextBox 13", "OUR INNOVATIVE SOLUTIONS")
    set_text(slide, "TextBox 11", "HCS-Verified AI Agents")
    set_text(slide, "TextBox 12",
        "Every trade signal is written to a Hedera HCS topic — immutable, "
        "timestamped, and publicly auditable. No central party can alter history."
    )
    set_text(slide, "TextBox 14", "NFT Strategy Marketplace")
    set_text(slide, "TextBox 15",
        "Trade strategies as Hedera HTS NFTs. Creators earn 5% royalty on every "
        "secondary sale. Buyers get verified, on-chain performance records."
    )
    set_text(slide, "TextBox 16", "On-Chain Audit Trail")
    set_text(slide, "TextBox 17",
        "All agent deployments, configurations (stored on HFS), and contract calls "
        "(HSCS) form a complete, tamper-proof audit trail visible to anyone."
    )


def fill_slide_6(slide):
    """Hedera Services"""
    set_text(slide, "TextBox 33", "HEDERA SERVICES WE USE")
    set_text(slide, "TextBox 31", "HCS — Hedera Consensus Service")
    set_text(slide, "TextBox 32",
        "Real-time trade signals and agent events are published to HCS topics. "
        "Immutable ordering via aBFT consensus. Powers the live signal feed."
    )
    set_text(slide, "TextBox 34", "HTS — Hedera Token Service")
    set_text(slide, "TextBox 35",
        "Agent strategies are minted as NFTs (HIP-412). Enables the marketplace, "
        "ownership transfers, and automatic 5% royalty splits in native HBAR."
    )
    set_text(slide, "TextBox 36", "HFS — Hedera File Service")
    set_text(slide, "TextBox 37",
        "Agent configuration JSON is stored on HFS for tamper-proof, "
        "decentralized config management — referenced by HCS and HSCS."
    )
    set_text(slide, "TextBox 38", "HSCS — Smart Contract Service")
    set_text(slide, "TextBox 39",
        "AgentRegistry.sol and MockDEX.sol deployed on Hedera EVM. "
        "Handles on-chain agent registration, trade execution, and settlement."
    )


def fill_slide_7(slide):
    """Market Size"""
    set_text(slide, "TextBox 15", "SIZE OF MARKET")
    set_text(slide, "TextBox 14",
        "The global algorithmic trading market is growing at 12.9% CAGR. "
        "Crypto retail trading alone surpassed $2.1T in 2024. "
        "AI-powered DeFi tools represent the fastest-growing segment."
    )
    set_text(slide, "TextBox 16", "Total Addressable Market (TAM):  $8.3 Trillion")
    set_text(slide, "TextBox 17", "Serviceable Addressable Market (SAM):  $2.1 Trillion")
    set_text(slide, "TextBox 18", "Serviceable Obtainable Market (SOM):  $42 Billion")


def fill_slide_8(slide):
    """Competitors"""
    set_text(slide, "TextBox 22", "DIRECT COMPETITORS")
    set_text(slide, "TextBox 23",
        "3Commas, Cryptohopper, centralized AI bots.\n"
        "Target same retail crypto traders.\n"
        "No on-chain verifiability, no HCS audit trail, no NFT marketplace.\n"
        "Rely on centralized servers — single point of failure and manipulation."
    )
    set_text(slide, "TextBox 24", "INDIRECT COMPETITORS")
    set_text(slide, "TextBox 25",
        "Traditional quant funds and Web2 analytics platforms (TradingView, etc.).\n"
        "Solve similar needs but require coding expertise.\n"
        "No blockchain transparency, no creator monetization, no HBAR-native settlement."
    )
    set_text(slide, "TextBox 27", "Arcane")


def fill_slide_9(slide):
    """Competitive Advantages"""
    set_text(slide, "TextBox 20", "KEY COMPETITIVE ADVANTAGES")
    set_text(slide, "TextBox 19", "HCS-Verified Execution")
    set_text(slide, "TextBox 18",
        "Every signal is recorded on Hedera HCS — immutable, timestamped, "
        "and publicly verifiable. Competitors offer zero on-chain proof."
    )
    set_text(slide, "TextBox 22", "AI Strategy Engine")
    set_text(slide, "TextBox 21",
        "LangGraph ReAct AI with Pyth Oracle feeds, EMA/RSI/MACD indicators, "
        "and Kelly Criterion risk sizing — institutional quality, retail accessible."
    )
    set_text(slide, "TextBox 24", "NFT Strategy Marketplace")
    set_text(slide, "TextBox 23",
        "Strategies trade as Hedera HTS NFTs. Creators earn 5% royalty perpetually. "
        "Buyers receive verified, on-chain performance history."
    )
    set_text(slide, "TextBox 26", "Zero-Code Deployment")
    set_text(slide, "TextBox 25",
        "Natural language agent creation powered by Google Gemini 2.5 Flash. "
        "Deploy a live trading agent in under 2 minutes — no coding required."
    )


def fill_slide_10(slide):
    """Traction"""
    set_text(slide, "TextBox 5", "TRACTION")
    set_text(slide, "TextBox 4",
        "Fully functional MVP deployed on Hedera Testnet.\n\n"
        "• HCS topics created and live — real-time signal feed operational\n"
        "• AgentRegistry.sol & MockDEX.sol deployed on Hedera EVM\n"
        "• NFT strategy marketplace: mint, list, buy, and royalty split working\n"
        "• WalletConnect v2 integration — wallet sign-in live\n"
        "• AI agent runs with Pyth Oracle + LangGraph ReAct confirmed\n"
        "• Candlestick chart analytics page live\n"
        "• End-to-end trade cycle verified: signal → HCS → execution → audit"
    )


def fill_slide_11(slide):
    """Revenue Model — heading only, visual content stays"""
    set_text(slide, "TextBox 13", "REVENUE MODEL")


def fill_slide_12(slide):
    """Roadmap"""
    set_text(slide, "TextBox 13", "ROADMAP")
    set_text(slide, "TextBox 11", "Q1 2025")
    set_text(slide, "TextBox 12",
        "Core platform architecture built. Hedera HCS integration, "
        "LangGraph AI agent framework, and smart contract foundation on testnet."
    )
    set_text(slide, "TextBox 14", "Q2 2025")
    set_text(slide, "TextBox 15",
        "Full MVP deployed. NFT marketplace, WalletConnect, candlestick analytics, "
        "HFS config storage, and end-to-end trade cycle verified."
    )
    set_text(slide, "TextBox 16", "Q3 2025  (Now)")
    set_text(slide, "TextBox 17",
        "Hedera Hackathon submission. Community beta launch, "
        "early adopter onboarding, and market feedback cycles initiated."
    )
    set_text(slide, "TextBox 18", "Q4 2025 — Mainnet")
    set_text(slide, "TextBox 19",
        "Hedera mainnet launch. SaucerSwap V2 DEX integration, "
        "paid subscription tiers, institutional API access, and HCS-10 OpenConvAI agents."
    )


def fill_slide_13(slide):
    """Use of Funds"""
    set_text(slide, "TextBox 5", "USE OF FUNDS")
    set_text(slide, "TextBox 4",
        "Funds raised will be allocated to accelerate growth and reach mainnet:\n\n"
        "• 40%  Product Development — AI model improvements, smart contract audits, "
        "SaucerSwap integration, mobile app\n"
        "• 30%  Go-To-Market — Community building, creator incentives, "
        "partnership with Hedera ecosystem projects\n"
        "• 20%  Infrastructure — Hedera node operations, Redis/database scaling, "
        "Pyth Oracle premium feeds\n"
        "• 10%  Legal & Compliance — Token legal structure, "
        "smart contract audit by third-party, regulatory guidance"
    )


def fill_slide_14(slide):
    """Meet the Team"""
    set_text(slide, "TextBox 14", "MEET THE TEAM")
    set_text(slide, "TextBox 11",
        "For questions, demos, or partnership opportunities, reach out directly."
    )
    # Primary member
    set_text(slide, "TextBox 12", "Founder & Full-Stack Developer")
    set_text(slide, "TextBox 13", "NISHANT GUPTA")
    # Clear unused slots
    set_text(slide, "TextBox 15", "")
    set_text(slide, "TextBox 16", "")
    set_text(slide, "TextBox 17", "")
    set_text(slide, "TextBox 18", "")
    set_text(slide, "TextBox 19", "")
    set_text(slide, "TextBox 20", "")


def fill_slide_15(slide):
    """Thank You"""
    set_text(slide, "TextBox 9",  "THANK YOU")
    set_text(slide, "TextBox 15", "for your time and attention")
    set_text(slide, "TextBox 10", "Arcane")
    set_text(slide, "TextBox 11", "Built by Nishant Gupta")
    set_text(slide, "TextBox 12", "nishantgupta1965@gmail.com")
    set_text(slide, "TextBox 13", "AI & Agents Track  |  Hedera Hackathon 2025")
    set_text(slide, "TextBox 14", "Powered by Hedera Hashgraph")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
FILLERS = [
    fill_slide_1,
    fill_slide_2,
    fill_slide_3,
    fill_slide_4,
    fill_slide_5,
    fill_slide_6,
    fill_slide_7,
    fill_slide_8,
    fill_slide_9,
    fill_slide_10,
    fill_slide_11,
    fill_slide_12,
    fill_slide_13,
    fill_slide_14,
    fill_slide_15,
]

if __name__ == "__main__":
    prs = Presentation(TEMPLATE)

    if len(prs.slides) != len(FILLERS):
        print(f"WARNING: template has {len(prs.slides)} slides, expected {len(FILLERS)}")

    for i, (slide, filler) in enumerate(zip(prs.slides, FILLERS)):
        filler(slide)
        print(f"  Filled slide {i+1}: {filler.__doc__}")

    prs.save(OUTPUT)
    print(f"\nSaved -> {OUTPUT}")
